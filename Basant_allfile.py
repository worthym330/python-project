from tkinter import *
from tkinter import filedialog, messagebox, scrolledtext
from PIL import Image, ImageTk # install pillow
import os
from stegano import lsb # install stegano
import cv2
import pandas as pd
from docx import Document
import PyPDF2
from pptx import Presentation
import pygame
from pygame import mixer
import wave
import numpy as np
# pip install pandas python-docx PyPDF2 python-pptx openpyxl pygame

root=Tk()
root.title("Steganography - Hide a Secret Message in a File")
root.geometry("700x500+150+180")
root.resizable(False,False)
root.configure(bg="#2f4155")

filename = None
hidden_data = None

def show_file():
    global filename
    filename = filedialog.askopenfilename(initialdir=os.getcwd(), title='Select File', filetypes=(("All files", "*.*"),))
    if filename:
        try:
            if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.mp4', '.avi', '.mov','.mp3', '.wav')):
                show_image_or_video()
            elif filename.lower().endswith(('.txt', '.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.csv')):
                show_text()
            else:
                messagebox.showinfo("Info", f"Selected file: {filename}")
        except Exception as e:
            messagebox.showerror("Error", f"Error opening file: {str(e)}")


def show_image_or_video():
    global filename
    f = Frame(root, bd=3, bg="black", width=340, height=280, relief="groove")
    f.place(x=10, y=80)

    if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
        img = Image.open(filename)
        img_width, img_height = img.size

        # Resize image to fit within frame without distortion
        aspect_ratio = img_width / img_height
        if img_width > img_height:
            new_width = 320
            new_height = int(new_width / aspect_ratio)
        else:
            new_height = 270
            new_width = int(new_height * aspect_ratio)
        img = img.resize((new_width, new_height))

        img = ImageTk.PhotoImage(img)
        lbl = Label(f, bg="black", image=img)
        lbl.place(x=(340 - new_width) // 2, y=(280 - new_height) // 2)
        lbl.image = img

    elif filename.lower().endswith(('.mp4', '.avi', '.mov','.mkv')):
        cap = cv2.VideoCapture(filename)
        ret, frame = cap.read()
        if ret:
            frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            img = Image.fromarray(frame)
            img_width, img_height = img.size
            
            # Resize image to fit within label without distortion
            aspect_ratio = img_width / img_height
            if img_width > img_height:
                new_width = 250
                new_height = int(new_width / aspect_ratio)
            else:
                new_height = 250
                new_width = int(new_height * aspect_ratio)
            img = img.resize((new_width, new_height))
            
            img = ImageTk.PhotoImage(img)
            lbl = Label(f, bg="black", image=img)
            lbl.configure(image=img, width=new_width, height=new_height)
            lbl.place = img
        cap.release()

    elif filename.lower().endswith(('.mp3', '.wav','m4a')):
        pygame.init()
        mixer.init()
        mixer.music.load(filename)
        mixer.music.play()

def show_text():
    global filename
    text_frame = Frame(root, bd=3, bg="#2f4155", width=340, height=280, relief="groove")
    text_frame.place(x=10, y=80)

    lbl = Text(text_frame, font="Roboto 12", bg="white", fg="black", relief="groove", wrap="word")
    lbl.place(x=10, y=10, width=320, height=260)

    if filename.lower().endswith(('.xls', '.xlsx', '.csv')):
        # Read Excel or CSV file and display content
        try:
            df = pd.read_excel(filename) if filename.lower().endswith(('.xls', '.xlsx')) else pd.read_csv(filename)
            text_content = df.to_string(index=False)
            lbl.insert("1.0", text_content)
        except Exception as e:
            messagebox.showerror("Error", f"Error reading file: {str(e)}")
    
    elif filename.lower().endswith(('.doc', '.docx')):
        # Read Word document and display content
        try:
            doc = Document(filename)
            text_content = '\n'.join([para.text for para in doc.paragraphs])
            lbl.insert("1.0", text_content)
        except Exception as e:
            messagebox.showerror("Error", f"Error reading file: {str(e)}")
    elif filename.lower().endswith('.pdf'):
        # Read PDF file and display content
        try:
            with open(filename, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                num_pages = len(pdf_reader.pages)
                text_content = ''
                for page_num in range(num_pages):
                    text_content += pdf_reader.pages[page_num].extract_text()
            lbl.insert("1.0", text_content)
        except Exception as e:
            messagebox.showerror("Error", f"Error reading file: {str(e)}")
    elif filename.lower().endswith('.txt'):
        # Read text file and display content
        try:
            with open(filename, 'r') as file:
                text_content = file.read()
            lbl.insert("1.0", text_content)
        except Exception as e:
            messagebox.showerror("Error", f"Error reading file: {str(e)}")
    elif filename.lower().endswith(('.ppt', '.pptx')):
        # Read PowerPoint presentation and display content
        try:
            prs = Presentation(filename)
            text_content = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text_content += shape.text + '\n'
            text1.insert("1.0", text_content)
        except Exception as e:
            messagebox.showerror("Error", f"Error reading file: {str(e)}")
    else:
        messagebox.showinfo("Info", f"Selected file: {filename}")


def hide_message_in_image():
    global hidden_data
    global filename
    message=text1.get(1.0, END)
    try:
        hidden_data = lsb.hide(str(filename), message)
        messagebox.showinfo("Success", "Data hidden successfully!")
    except Exception as e:
        messagebox.showerror("Error", f"Error hiding data: {str(e)}")

def hide_message_in_audio():
    global filename
    global hidden_data
    try:
        # Open the audio file
        audio = wave.open(filename, 'rb')

        # Read audio frames
        frames = audio.readframes(audio.getnframes())

        # Convert frames to numpy array
        audio_array = np.frombuffer(frames, dtype=np.int16)

        # Convert message to binary
        binary_message = ''.join(format(ord(char), '08b') for char in hidden_data)

        # Ensure the message fits in the audio file
        if len(binary_message) > len(audio_array):
            raise ValueError("Message too large to hide in audio file.")

        # Hide the message in the least significant bits
        for i in range(len(binary_message)):
            audio_array[i] = (audio_array[i] & ~1) | int(binary_message[i])

        # Convert the numpy array back to bytes
        modified_frames = audio_array.tobytes()

        # Create a new wave file with the hidden message
        with wave.open('hidden_audio.wav', 'wb') as output_audio:
            output_audio.setparams(audio.getparams())
            output_audio.writeframes(modified_frames)
        
        print("Message hidden successfully in audio.")
    except Exception as e:
        print("Error hiding message in audio:", str(e))

def hide_message_in_file():
    global filename
    global hidden_data
    try:
        with open(filename, "a") as file:
            file.write("\n" + hidden_data)
        print("Message hidden successfully in file.")
    except Exception as e:
        print("Error hiding message in file:", str(e))

def hide():
    global filename
    global message
    try:
        if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
            hide_message_in_image()
        elif filename.lower().endswith(('.wav', '.m4a')):
            hide_message_in_audio()
        elif filename.lower().endswith(('.txt', '.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.csv')):
            hide_message_in_file()
        elif filename.lower().endswith(('.mkv', '.mov', '.mp4')):
            hide_message_in_video()
        else:
            print("Unsupported file format!")
    except Exception as e:
        print("Error hiding message:", str(e))

def reveal_message_from_image():
    global filename
    try:
        clear_message = lsb.reveal(filename)
        text1.delete(1.0, END)
        text1.insert(END, clear_message)
    except Exception as e:
        messagebox.showerror("Error", f"Error revealing data: {str(e)}")

def reveal_message_from_audio():
    global filename
    try:
        # Open the audio file
        audio = wave.open(filename, 'rb')

        # Read audio frames
        frames = audio.readframes(audio.getnframes())

        # Convert frames to numpy array
        audio_array = np.frombuffer(frames, dtype=np.int16)

        binary_message = ''
        for sample in audio_array:
            binary_message += str(sample & 1)

        message = ''
        for i in range(0, len(binary_message), 8):
            byte = binary_message[i:i+8]
            message += chr(int(byte, 2))

        print("Hidden message in audio:", message)
    except Exception as e:
        print("Error revealing message from audio:", str(e))

def reveal_message_from_file():
    global filename
    try:
        with open(filename, "r") as file:
            lines = file.readlines()
            print("Hidden message in file:", "".join(lines))
    except Exception as e:
        print("Error revealing message from file:", str(e))

def show():
    global filename
    try:
        if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
            reveal_message_from_image()
        elif filename.lower().endswith(('.wav', '.m4a')):
            reveal_message_from_audio()
        elif filename.lower().endswith(('.txt', '.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.csv')):
            reveal_message_from_file()
        elif filename.lower().endswith(('.mkv', '.mov', '.mp4')):
            hide_message_in_video()
        else:
            print("Unsupported file format!")
    except Exception as e:
        print("Error revealing message:", str(e))


def save_image():
    global hidden_data
    if hidden_data:
        try:
            output_filename = filedialog.asksaveasfilename(defaultextension=".png", filetypes=(("PNG files", "*.png"), ("All files", "*.*")))
            if output_filename:
                hidden_data.save(output_filename)
                messagebox.showinfo("Success", "Image saved successfully!")
        except Exception as e:
            messagebox.showerror("Error", f"Error saving image: {str(e)}")
    else:
        messagebox.showerror("Error", "No hidden data to save!")

def save_audio():
    global hidden_data
    if hidden_data:
        try:
            output_filename = filedialog.asksaveasfilename(defaultextension=".wav", filetypes=(("WAV files", "*.wav"), ("All files", "*.*")))
            if output_filename:
                with wave.open(output_filename, 'wb') as output_audio:
                    output_audio.setparams(hidden_data.getparams())
                    output_audio.writeframes(hidden_data.readframes(hidden_data.getnframes()))
                messagebox.showinfo("Success", "Audio saved successfully!")
        except Exception as e:
            messagebox.showerror("Error", f"Error saving audio: {str(e)}")
    else:
        messagebox.showerror("Error", "No hidden data to save!")

def save_text():
    global hidden_data
    if hidden_data:
        try:
            output_filename = filedialog.asksaveasfilename(defaultextension=".txt", filetypes=(("Text files", "*.txt"), ("All files", "*.*")))
            if output_filename:
                with open(output_filename, "w") as file:
                    file.write(hidden_data)
                messagebox.showinfo("Success", "Text saved successfully!")
        except Exception as e:
            messagebox.showerror("Error", f"Error saving text: {str(e)}")
    else:
        messagebox.showerror("Error", "No hidden data to save!")

def save_video():
    # Implement saving video data here
    pass

def save():
    global filename
    if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
        save_image()
    elif filename.lower().endswith('.wav'):
        save_audio()
    elif filename.lower().endswith('.txt'):
        save_text()
    elif filename.lower().endswith(('.mp4', '.avi', '.mov')):
        save_video()
    else:
        messagebox.showinfo("Info", "Unsupported file format for saving.")


#icon
image_icon=PhotoImage(file="E:\VS code\Projects\project data for stegano\logo.jpg")
root.iconphoto(False,image_icon)


#logo
logo=PhotoImage(file="E:\VS code\Projects\project data for stegano\logo.png")
Label(root,image=logo,bg="#2f4155").place(x=10,y=0)
Label(root,text="CYBER SCIENCE",bg="#2f4155",fg="white",font="arial 25 bold").place(x=100,y=20)


# First Frame
f=Frame(root,bd=3,bg="black",width=340,height=280,relief=GROOVE)
f.place(x=10,y=80)

lbl=Label(f,bg="black")
lbl.place(x=40,y=10)

# Second Frame
frame2=Frame(root,bd=3,bg="black",width=340,height=280,relief=GROOVE)
frame2.place(x=350,y=80)

text1=Text(frame2,font="Roboto 20",bg="white",fg="black",relief=GROOVE,wrap=WORD)
text1.place(x=0,y=0,width=320,height=295)

scrollbar1=Scrollbar(frame2)
scrollbar1.place(x=320,y=0,height=300)

scrollbar1.configure(command=text1.yview)
text1.configure(yscrollcommand=scrollbar1.set)

# Third Frame
frame3=Frame(root,bd=3,bg="#2f4155",width=330,height=100,relief=GROOVE)
frame3.place(x=10,y=370)

Button(frame3,text="Open File",width=10,height=2,font="arial 14 bold",command=show_file).place(x=20,y=30)
Button(frame3,text="Save File",width=10,height=2,font="arial 14 bold",command=save).place(x=180,y=30)
Label(frame3,text="Select Image, Video, Audio, or Text File",bg="#2f4155",fg="yellow").place(x=20,y=5)

# Fourth Frame
frame4=Frame(root,bd=3,bg="#2f4155",width=340,height=280,relief=GROOVE)
frame4.place(x=360,y=370)

Button(frame4,text="Hide Data",width=10,height=2,font="arial 14 bold",command=hide).place(x=20,y=30)
Button(frame4,text="Show Data",width=10,height=2,font="arial 14 bold",command=show).place(x=180,y=30)
Label(frame4,text="Hide and Show Data",bg="#2f4155",fg="yellow").place(x=20,y=5)

root.mainloop()